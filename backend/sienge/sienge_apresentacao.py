from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import plotly.io as pio
import pandas as pd
import tempfile
import os

from sienge.sienge_ia import gerar_analise_financeira


# ============================================================
# 🧩 GERAÇÃO DE APRESENTAÇÃO POWERPOINT
# ============================================================
def gerar_apresentacao_ppt(df: pd.DataFrame, resumo_dre: dict):
    """
    Gera uma apresentação PowerPoint (PPTX) com gráficos e textos da IA.
    Retorna bytes do arquivo para download ou exportação.
    """

    prs = Presentation()

    # === SLIDE 1 — TÍTULO ===
    slide_title = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide_title.shapes.title
    subtitle = slide_title.placeholders[1]
    title.text = "Relatório Financeiro — Constru.IA"
    subtitle.text = "Gerado automaticamente via API Sienge e IA"

    # === SLIDE 2 — DRE ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    title.text = "Resumo Financeiro (DRE)"
    content = slide2.placeholders[1]
    content.text = (
        f"💵 Receitas: {resumo_dre.get('receitas', 'R$ 0,00')}\n"
        f"💸 Despesas: {resumo_dre.get('despesas', 'R$ 0,00')}\n"
        f"📈 Lucro: {resumo_dre.get('lucro', 'R$ 0,00')}"
    )

    # === SLIDES DE GRÁFICOS ===
    def add_chart_slide(df_group, title_text, group_col, filename="grafico.png"):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = title_text

        # cria o gráfico com Plotly
        fig = pio.to_image(
            pio.Figure().add_bar(x=df_group[group_col], y=df_group["valor_total"]),
            format="png", width=1000, height=500
        )
        tmp_path = os.path.join(tempfile.gettempdir(), filename)
        with open(tmp_path, "wb") as f:
            f.write(fig)

        slide.shapes.add_picture(tmp_path, Inches(1), Inches(2), width=Inches(8))
        os.remove(tmp_path)

        analise = gerar_analise_financeira(title_text, df_group)
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(6), Inches(8.5), Inches(2))
        tf = tx_box.text_frame
        p = tf.add_paragraph()
        p.text = analise
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT

    # Gera slides por categoria
    if "obra" in df.columns:
        obra_data = df.groupby("obra")["valor_total"].sum().reset_index()
        add_chart_slide(obra_data, "Gastos por Obra", "obra")

    if "centro_custo" in df.columns:
        cc_data = df.groupby("centro_custo")["valor_total"].sum().reset_index()
        add_chart_slide(cc_data, "Gastos por Centro de Custo", "centro_custo")

    if "fornecedor" in df.columns:
        forn_data = df.groupby("fornecedor")["valor_total"].sum().reset_index()
        add_chart_slide(forn_data, "Top Fornecedores", "fornecedor")

    # === SALVA EM MEMÓRIA ===
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
