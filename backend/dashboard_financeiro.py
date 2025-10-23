import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

def gerar_apresentacao_ppt(df: pd.DataFrame, dre: dict) -> bytes:
    """Gera um arquivo PowerPoint com gráficos e KPIs financeiros."""
    try:
        prs = Presentation()
        slide_title = prs.slide_layouts[0]
        slide_content = prs.slide_layouts[5]

        # --- Slide 1: Título ---
        slide = prs.slides.add_slide(slide_title)
        slide.shapes.title.text = "Relatório Financeiro — Constru.IA"
        slide.placeholders[1].text = "Resumo Automático via API Sienge + IA"

        # --- Slide 2: Indicadores Financeiros ---
        slide = prs.slides.add_slide(slide_content)
        slide.shapes.title.text = "📊 Indicadores Gerais"
        tf = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3)).text_frame
        tf.text = f"💵 Receitas: {dre.get('receitas', 'R$ 0,00')}\n💸 Despesas: {dre.get('despesas', 'R$ 0,00')}\n📈 Lucro: {dre.get('lucro', 'R$ 0,00')}"

        # --- Slide 3: Gráfico de Gastos por Categoria ---
        if "centro_custo" in df.columns:
            data = df.groupby("centro_custo")["valor_total"].sum().reset_index().sort_values("valor_total", ascending=False).head(8)
            chart_data = CategoryChartData()
            chart_data.categories = data["centro_custo"]
            chart_data.add_series("Gastos", data["valor_total"])

            slide = prs.slides.add_slide(slide_content)
            slide.shapes.title.text = "🏗️ Gastos por Centro de Custo"
            x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
            graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
            graphic_frame.has_legend = True

        # --- Slide 4: Tabela Resumo ---
        slide = prs.slides.add_slide(slide_content)
        slide.shapes.title.text = "📋 Amostra de Dados"
        top = Inches(1.5)
        left = Inches(0.8)
        width = Inches(8.5)
        height = Inches(0.8)
        table_data = df.head(8)
        rows, cols = table_data.shape
        table = slide.shapes.add_table(rows+1, cols, left, top, width, height).table
        for j, col in enumerate(table_data.columns):
            table.cell(0, j).text = col
        for i, row in table_data.iterrows():
            for j, value in enumerate(row):
                table.cell(i+1, j).text = str(value)[:25]

        # --- Slide 5: Conclusão ---
        slide = prs.slides.add_slide(slide_content)
        slide.shapes.title.text = "💡 Conclusões e Recomendações"
        slide.placeholders[0].text = (
            "• Otimizar gastos nas obras de maior peso.\n"
            "• Reavaliar fornecedores com concentração de custos.\n"
            "• Priorizar obras com maior retorno e menor inadimplência.\n"
            "• Reduzir despesas administrativas recorrentes."
        )

        # Exporta o arquivo em bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    except Exception as e:
        print("Erro ao gerar PPT:", e)
        return None
